from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ===== SLIDE 1: TITLE =====
s1 = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide layout
s1.shapes.title.text = "AI MENTAL HEALTH UK"
s1.placeholders[1].text = "Mental Health Risk Prediction Model\nAIML-4 Demo Walkthrough\nMay 2026\n\nTeam: Dinisha Jain | Mohit"

# ===== SLIDE 2: PROJECT OVERVIEW =====
s2 = prs.slides.add_slide(prs.slide_layouts[1])
s2.shapes.title.text = "Project Overview"
tf = s2.placeholders[1].text_frame
p = tf.paragraphs[0]
p.text = "Purpose: Predict 30-day readmission risk for mental health patients"
p.font.size = Pt(18)
for text in ["Model: Ensemble (XGBoost + RF + LogReg)", "Features: 14 input features", "Target: F1 > 0.80", "Data: NHS MHSDS"]:
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(18)

# ===== SLIDE 3: WEEK 1 =====
s3 = prs.slides.add_slide(prs.slide_layouts[1])
s3.shapes.title.text = "Week 1: Data Collection & EDA"
tf = s3.placeholders[1].text_frame
p = tf.paragraphs[0]
p.text = "AIML-1: Data Analysis"
p.font.size = Pt(18)
for text in ["496 patient records analyzed", "PHQ-9: Mean=12, SD=7", "GAD-7: Mean=10, SD=6", "14 features engineered"]:
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(18)

# ===== SLIDE 4: WEEK 1 IMAGE =====
s4 = prs.slides.add_slide(prs.slide_layouts[6])
title4 = s4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
title4.text = "Week 1: Feature Analysis"
title4.text_frame.paragraphs[0].font.size = Pt(32)
title4.text_frame.paragraphs[0].font.bold = True
try:
    s4.shapes.add_picture('Week 1/reports/feature_target_bars.png', Inches(1), Inches(1.5), width=Inches(11))
except:
    txt = s4.shapes.add_textbox(Inches(2), Inches(3), Inches(9), Inches(1))
    txt.text = "Image: feature_target_bars.png"

# ===== SLIDE 5: WEEK 2 =====
s5 = prs.slides.add_slide(prs.slide_layouts[1])
s5.shapes.title.text = "Week 2: Baseline Models"
tf = s5.placeholders[1].text_frame
p = tf.paragraphs[0]
p.text = "AIML-2: Model Development"
p.font.size = Pt(18)
for text in ["Logistic Regression: F1=0.72, AUC=0.78", "Random Forest: F1=0.75, AUC=0.80", "XGBoost: F1=0.78, AUC=0.82"]:
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(18)

# ===== SLIDE 6: ROC CURVE =====
s6 = prs.slides.add_slide(prs.slide_layouts[6])
title6 = s6.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
title6.text = "Week 2: ROC Curve"
title6.text_frame.paragraphs[0].font.size = Pt(32)
title6.text_frame.paragraphs[0].font.bold = True
try:
    s6.shapes.add_picture('reports/xgboost_roc_curve.png', Inches(2), Inches(1.5), width=Inches(9))
except:
    pass

# ===== SLIDE 7: WEEK 3 =====
s7 = prs.slides.add_slide(prs.slide_layouts[1])
s7.shapes.title.text = "Week 3: Ensemble + Tuning"
tf = s7.placeholders[1].text_frame
p = tf.paragraphs[0]
p.text = "AIML-2: Ensemble Model"
p.font.size = Pt(18)
for text in ["Voting Ensemble: XGBoost + RF + LogReg", "20 Optuna trials", "Training F1=1.00, AUC=1.00", "Saved: ensemble_v3.pkl"]:
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(18)

# ===== SLIDE 8: ENSEMBLE IMAGE =====
s8 = prs.slides.add_slide(prs.slide_layouts[6])
title8 = s8.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
title8.text = "Week 3: Ensemble Comparison"
title8.text_frame.paragraphs[0].font.size = Pt(32)
title8.text_frame.paragraphs[0].font.bold = True
try:
    s8.shapes.add_picture('reports/ensemble_comparison.png', Inches(2), Inches(1.5), width=Inches(9))
except:
    pass

# ===== SLIDE 9: SHAP =====
s9 = prs.slides.add_slide(prs.slide_layouts[6])
title9 = s9.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
title9.text = "Week 3: SHAP Analysis"
title9.text_frame.paragraphs[0].font.size = Pt(32)
title9.text_frame.paragraphs[0].font.bold = True
try:
    s9.shapes.add_picture('reports/shap_bar.png', Inches(2), Inches(1.5), width=Inches(9))
except:
    pass

# ===== SLIDE 10: WEEK 4 VALIDATION =====
s10 = prs.slides.add_slide(prs.slide_layouts[1])
s10.shapes.title.text = "Week 4: Model Validation"
tf = s10.placeholders[1].text_frame
p = tf.paragraphs[0]
p.text = "AIML-1: Validation"
p.font.size = Pt(18)
for text in ["MHSDS Validation: F1=0.34, AUC=0.51", "Model Card created", "Load testing completed"]:
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(18)

# ===== SLIDE 11: BIAS =====
s11 = prs.slides.add_slide(prs.slide_layouts[6])
title11 = s11.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
title11.text = "Week 4: Bias Audit"
title11.text_frame.paragraphs[0].font.size = Pt(32)
title11.text_frame.paragraphs[0].font.bold = True
try:
    s11.shapes.add_picture('reports/bias_audit_plot.png', Inches(2), Inches(1.5), width=Inches(9))
except:
    pass

# ===== SLIDE 12: DEMO =====
s12 = prs.slides.add_slide(prs.slide_layouts[1])
s12.shapes.title.text = "AIML-4: Live Demo"
tf = s12.placeholders[1].text_frame
p = tf.paragraphs[0]
p.text = "Run: python Week 4/demo_notebook.py"
p.font.size = Pt(18)
for text in ["P001: PHQ-9=18, GAD-7=15 → Risk 9.95 (HIGH)", "P002: PHQ-9=3, GAD-7=2 → Risk 0.02 (LOW)", "P003: PHQ-9=20, GAD-7=18 → Risk 9.96 (HIGH)"]:
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(18)

# ===== SLIDE 13: LIMITATIONS =====
s13 = prs.slides.add_slide(prs.slide_layouts[1])
s13.shapes.title.text = "Limitations"
tf = s13.placeholders[1].text_frame
p = tf.paragraphs[0]
p.text = "NOT clinically validated"
for text in ["Trained on synthetic data", "Validation metrics need improvement", "NHS approval pending", "NOT clinical replacement"]:
    p = tf.add_paragraph()
    p.text = text

# ===== SLIDE 14: THANK YOU =====
s14 = prs.slides.add_slide(prs.slide_layouts[0])
s14.shapes.title.text = "Thank You!"
s14.placeholders[1].text = "Dinisha Jain\nMohit\nMay 2026\n\nQuestions?"

prs.save("AI_Mental_Health_Presentation.pptx")
print("Created!")